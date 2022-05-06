import os

from pptx import Presentation
from pptx.dml.color import RGBColor

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)

file_path = os.getenv("FILE_PATH")

if file_path is None: file_path = input("Enter file path to darken >> ")
presentation = Presentation(file_path)

for slide in presentation.slides:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BLACK
    for shape in slide.shapes:
        # TODO: Change outline of shape to #WHITE
        if not shape.has_text_frame: continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                font = run.font
                font.color.rgb = WHITE

folder = os.path.abspath(os.path.join(file_path, os.pardir))
file_name = os.path.basename(file_path)
new_file_path = f"{folder}{os.path.sep}dark.{file_name}"
presentation.save(new_file_path)
print(f"Successfully saved as: {new_file_path}")
